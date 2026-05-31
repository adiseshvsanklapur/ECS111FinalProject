"""Generate the ECS 111 final-project slide deck as a .pptx file.

This builds ``slides/ECS111_slides.pptx`` (about 11 slides) for the project
"Chain-of-Thought Prompting vs. Fine-Tuning for Table Reasoning in Small
Language Models". The generator is idempotent: re-running overwrites the
deck cleanly. Run with ``python scripts/make_slides.py``.

No numeric results are baked in. Every number is written as a bracketed
placeholder token so the team can inject real values from a Colab T4 run
(2 seeds, 1000 eval, 8000 train) later. The placeholder convention is a
square-bracketed label of the form ``[METRIC: condition model]``.

Placeholder token reference (token -> result it maps to)
--------------------------------------------------------
WikiTableQuestions, Exact Match (mean +/- std over 2 seeds):
  [EM: baseline base]          EM for C0 baseline, Flan-T5-base
  [EM: baseline large]         EM for C0 baseline, Flan-T5-large
  [EM: cot_plain base]         EM for CoT plain prompting, base
  [EM: cot_plain large]        EM for CoT plain prompting, large
  [EM: cot_structured base]    EM for CoT structured prompting, base
  [EM: cot_structured large]   EM for CoT structured prompting, large
  [EM: finetune_answers base]  EM for fine-tune answers-only, base
  [EM: finetune_traces base]   EM for fine-tune with reasoning traces, base

WikiTableQuestions, token-level F1 (mean +/- std over 2 seeds):
  [F1: baseline base]          token-F1 for C0 baseline, base
  [F1: baseline large]         token-F1 for C0 baseline, large
  [F1: cot_plain base]         token-F1 for CoT plain, base
  [F1: cot_plain large]        token-F1 for CoT plain, large
  [F1: cot_structured base]    token-F1 for CoT structured, base
  [F1: cot_structured large]   token-F1 for CoT structured, large
  [F1: finetune_answers base]  token-F1 for fine-tune answers-only, base
  [F1: finetune_traces base]   token-F1 for fine-tune with traces, base

Best-vs-baseline gap (success criterion 1: beat baseline by >= 5 EM points):
  [EM_GAP: best vs baseline]   best condition EM minus baseline EM, base

TabFact generalization, classification accuracy (zero TabFact training).
Only the two fine-tuned base models transfer to TabFact; baseline base is an
untrained floor. Large is prompt-only (fine-tuning it OOMs on a T4):
  [ACC: tabfact baseline base]         TabFact acc, untrained base floor
  [ACC: tabfact finetune_answers base] TabFact acc, fine-tune answers-only, base
  [ACC: tabfact finetune_traces base]  TabFact acc, fine-tune with traces, base
  [ACC: tabfact best]                  best TabFact acc across conditions (vs 60% criterion)

Error-type breakdown (share of errors by category, headline condition):
  [ERR: lookup]        share of errors that are lookup failures
  [ERR: aggregation]   share of errors that are aggregation failures
  [ERR: multihop]      share of errors that are multi-hop failures

Statistical tests:
  [MCNEMAR: cot vs finetune]   McNemar p-value, CoT-best vs fine-tune-best (base, WTQ)
  [KAPPA: chain quality]       Cohen's kappa, two raters on 100 sampled CoT chains

Compute:
  [TIME: per condition]   wall-clock per condition on Colab T4 (target < 2 hours)
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from src import analysis, report_fill  # noqa: E402

SLIDES_DIR = Path(__file__).resolve().parent.parent / "slides"
OUT_PATH = SLIDES_DIR / "ECS111_slides.pptx"

# Visual constants.
ACCENT = RGBColor(0x1F, 0x3A, 0x5F)   # deep blue for titles
SUBTLE = RGBColor(0x55, 0x55, 0x55)   # grey for sub-text
BODY = RGBColor(0x22, 0x22, 0x22)

TITLE = "Chain-of-Thought Prompting vs. Fine-Tuning for Table Reasoning in Small Language Models"
TEAM = ["Adisesh", "Amar", "Nikhil", "Sanjay", "Anant"]


def _set_run(run, size, bold=False, color=BODY):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title text box.
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = TITLE
    _set_run(r, 30, bold=True, color=ACCENT)

    # Course line.
    cb = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(8.8), Inches(0.6))
    cp = cb.text_frame.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = "ECS 111 — Final Project"
    _set_run(cr, 20, bold=True, color=BODY)

    # Team line.
    mb = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(8.8), Inches(1.0))
    mtf = mb.text_frame
    mtf.word_wrap = True
    mp = mtf.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    mr.text = "Team: " + ", ".join(TEAM)
    _set_run(mr, 18, color=SUBTLE)


def add_content_slide(prs, title, bullets):
    """Add a title + bullet slide. ``bullets`` is a list of (text, level)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title.
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.9))
    tp = tb.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    _set_run(tr, 28, bold=True, color=ACCENT)

    # Body.
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(8.9), Inches(5.3))
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for text, level in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        prefix = "• " if level == 0 else "– "
        r = p.add_run()
        r.text = prefix + text
        _set_run(r, 18 if level == 0 else 16, color=BODY)
        p.space_after = Pt(6)
    return slide


def add_table_slide(prs, title, headers, rows, note=None):
    """Add a title + table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.9))
    tp = tb.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    _set_run(tr, 28, bold=True, color=ACCENT)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    top = Inches(1.5)
    height = Inches(min(4.6, 0.45 * n_rows))
    gtable = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), top, Inches(9.0), height
    ).table

    for j, h in enumerate(headers):
        cell = gtable.cell(0, j)
        cell.text = h
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = gtable.cell(i, j)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = BODY

    if note:
        nb = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(9.0), Inches(0.8))
        ntf = nb.text_frame
        ntf.word_wrap = True
        np_ = ntf.paragraphs[0]
        nr = np_.add_run()
        nr.text = note
        _set_run(nr, 13, color=SUBTLE)
    return slide


def build(prs):
    # 1. Title.
    add_title_slide(prs)

    # 2. Motivation / problem.
    add_content_slide(
        prs,
        "Motivation: Table QA Is Hard for Small Models",
        [
            ("Tables are the default format for real data: medical records, finance, government data.", 0),
            ("Answering a table question is a multi-step pipeline that must all succeed.", 0),
            ("Understand the question, locate the right cells, pick the operation, format the answer.", 1),
            ("A failure at any step yields a confident wrong answer.", 1),
            ("Most published evidence for prompting and fine-tuning uses very large, expensive models.", 0),
            ("We study the under-explored regime: small models that run for free on Colab.", 0),
        ],
    )

    # 3. Research question.
    add_content_slide(
        prs,
        "Research Question",
        [
            ("When you cannot afford a 100B-parameter model, which lever helps table reasoning more: "
             "chain-of-thought prompting or supervised fine-tuning?", 0),
            ("And just as important, where does each one break?", 0),
            ("CoT risk: small models may copy surface patterns without the underlying logic.", 1),
            ("Fine-tuning risk: overfitting to dataset-specific shortcuts that collapse on new formats.", 1),
            ("Secondary theme: interpretability. CoT shows a checkable reasoning chain; "
             "fine-tuning returns only a final answer.", 0),
        ],
    )

    # 4. Datasets.
    add_table_slide(
        prs,
        "Datasets",
        ["Dataset", "Task", "Role", "What it tests"],
        [
            ["WikiTableQuestions", "Answer a question from a table",
             "Primary (train + eval)", "Filter, sort, compare, multi-hop reasoning"],
            ["TabFact", "True/false fact verification over a table",
             "Out-of-distribution test", "Generalization with zero TabFact training"],
        ],
        note="Tables are serialized to text (headers first, then rows) to fit the small models' "
             "token limits. Both load from HuggingFace datasets.",
    )

    # 5. Methods: the 6 conditions.
    add_table_slide(
        prs,
        "Methods: Six Conditions",
        ["Condition", "Model(s)", "Training", "Key detail"],
        [
            ["C0 Baseline", "base + large", "No", "Question + table only, no examples; sets the floor"],
            ["CoT plain", "base + large", "No", "Hand-written plain-English reasoning exemplars prepended"],
            ["CoT structured", "base + large", "No", "Same idea, structured step-by-step chain format"],
            ["Fine-tune answers-only", "base", "Yes", "Target is the final answer only"],
            ["Fine-tune + traces", "base", "Yes", "Target is a rule-based reasoning chain, then the answer"],
            ["Generalization", "above models", "No", "Evaluate on TabFact with zero TabFact training"],
        ],
        note="Fine-tuning is base only; Flan-T5-large runs out of memory on a free Colab T4.",
    )

    # 6. Experimental setup.
    add_content_slide(
        prs,
        "Experimental Setup",
        [
            ("Models: Flan-T5-base (250M), fine-tuned and prompted; Flan-T5-large (780M), prompting only.", 0),
            ("Decoding: greedy, temperature 0, for reproducibility.", 0),
            ("Seeds: 2 fixed seeds per condition; results reported as mean and standard deviation.", 0),
            ("Eval set: 1000 examples. Train set: 8000 examples (WTQ only; no TabFact training).", 0),
            ("Fine-tune hyperparameters: AdamW, learning rate 3e-4, batch 8 with grad-accum 4 "
             "(effective 32), 3 epochs.", 0),
            ("Hardware: free Google Colab T4 GPU, target under 2 hours per condition "
             "([TIME: per condition]).", 0),
        ],
    )

    # 7. Results: WTQ.
    add_table_slide(
        prs,
        "Results: WikiTableQuestions",
        ["Condition", "EM base", "EM large", "F1 base", "F1 large"],
        [
            ["Baseline", "[EM: baseline base]", "[EM: baseline large]",
             "[F1: baseline base]", "[F1: baseline large]"],
            ["CoT plain", "[EM: cot_plain base]", "[EM: cot_plain large]",
             "[F1: cot_plain base]", "[F1: cot_plain large]"],
            ["CoT structured", "[EM: cot_structured base]", "[EM: cot_structured large]",
             "[F1: cot_structured base]", "[F1: cot_structured large]"],
            ["Fine-tune answers", "[EM: finetune_answers base]", "n/a (OOM)",
             "[F1: finetune_answers base]", "n/a (OOM)"],
            ["Fine-tune traces", "[EM: finetune_traces base]", "n/a (OOM)",
             "[F1: finetune_traces base]", "n/a (OOM)"],
        ],
        note="Exact Match and token-F1, mean over 2 seeds. Success criterion 1: best condition beats "
             "baseline by at least 5 EM points ([EM_GAP: best vs baseline]).",
    )

    # 8. Results: TabFact generalization.
    add_content_slide(
        prs,
        "Results: Generalization on TabFact",
        [
            ("Classification accuracy on TabFact with zero TabFact training.", 0),
            ("We transfer the two fine-tuned base models to TabFact. Large is prompt-only "
             "(fine-tuning it OOMs on a free T4), so it has no transfer checkpoint.", 1),
            ("Baseline base, untrained floor: [ACC: tabfact baseline base].", 0),
            ("Fine-tune answers: [ACC: tabfact finetune_answers base].", 0),
            ("Fine-tune traces: [ACC: tabfact finetune_traces base].", 0),
            ("Success criterion 2: does any condition reach at least 60% on TabFact? "
             "Best is [ACC: tabfact best].", 0),
        ],
    )

    # 9. Error analysis.
    add_content_slide(
        prs,
        "Error Analysis",
        [
            ("Every wrong answer is labeled by failure type.", 0),
            ("Lookup: failed to find the right cell. Share of errors: [ERR: lookup].", 1),
            ("Aggregation: failed at sum, count, max, or similar. Share: [ERR: aggregation].", 1),
            ("Multi-hop: failed to chain multiple steps. Share: [ERR: multihop].", 1),
            ("Headline comparison validated with McNemar's test on paired predictions: "
             "p = [MCNEMAR: cot vs finetune].", 0),
            ("Reasoning-chain quality: 100 sampled chains, two raters, 0/1/2 scale; "
             "agreement Cohen's kappa = [KAPPA: chain quality].", 0),
        ],
    )

    # 10. Limitations & takeaways.
    add_content_slide(
        prs,
        "Limitations and Takeaways",
        [
            ("Fine-tuning was limited to the base model; large ran out of memory on a free T4.", 0),
            ("Two seeds give a rough variance estimate, not a tight confidence interval.", 0),
            ("Table serialization to text loses some structure compared to native table encoders.", 0),
            ("Rule-based reasoning traces cover unambiguous derivations, not every question type.", 0),
            ("Takeaway: at this model size, prompting and fine-tuning trade off differently between "
             "in-distribution accuracy and out-of-distribution generalization.", 0),
        ],
    )

    # 11. Conclusion + contributions.
    add_content_slide(
        prs,
        "Conclusion and Contributions",
        [
            ("A controlled comparison of CoT prompting and fine-tuning at 250M and 780M parameters "
             "under strict free-Colab limits.", 0),
            ("A cross-dataset generalization test: train on WTQ, evaluate on TabFact with zero "
             "TabFact training.", 0),
            ("An error-type breakdown (lookup, aggregation, multi-hop) plus McNemar and Cohen's kappa.", 0),
            ("Fully reproducible: fixed seeds, greedy decoding, and notebooks that run end to end.", 0),
            ("Team: " + ", ".join(TEAM) + ".", 0),
        ],
    )


def load_token_map():
    """Token -> real value from results/*.json, or None if there are no results yet."""
    results = analysis.load_results()
    if not results:
        return None
    cq_path = Path("results/chain_quality.json")
    chain_quality = (json.loads(cq_path.read_text()) if cq_path.exists()
                     else {"kappa": 0.0, "mean_a": 0.0, "mean_b": 0.0})
    return report_fill.build_token_map(results, chain_quality)


def _sub_text_frame(text_frame, token_map):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            for token, value in token_map.items():
                if token in run.text:
                    run.text = run.text.replace(token, value)


def substitute_tokens(prs, token_map):
    """Replace bracket tokens in every text box and table cell with real values."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _sub_text_frame(shape.text_frame, token_map)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _sub_text_frame(cell.text_frame, token_map)


def main():
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    build(prs)
    token_map = load_token_map()
    if token_map:
        substitute_tokens(prs, token_map)
    else:
        print("No results yet -- deck built with placeholder tokens.")
    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH} with {len(prs.slides._sldIdLst)} slides.")


if __name__ == "__main__":
    main()
